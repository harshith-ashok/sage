"""MS Office support, part 2: a real .pptx slide-deck writer (python-pptx) —
the other half of "python-docx/pptx/openpyxl" named in this project's Stack
from the start. A title slide plus one slide per {heading, bullets} —
deliberately simple (no themes/images/charts) to match how the docx/xlsx
writers stay plain: a real, openable Office file with the actual content
structured onto real slide placeholders, not a decorative deck.
"""

import os

from pptx import Presentation
from pptx.util import Pt


def write_pptx(path: str, title: str, subtitle: str, slides: list[dict]) -> None:
    """`slides` is a list of {"heading": str, "bullets": list[str]} dicts,
    one per content slide after the title slide."""
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    if len(title_slide.placeholders) > 1:
        title_slide.placeholders[1].text = subtitle or ""

    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # "Title and Content"
        slide.shapes.title.text = str(slide_data.get("heading", ""))
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = slide_data.get("bullets") or []
        for i, bullet in enumerate(bullets):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = str(bullet)
            p.font.size = Pt(18)

    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
