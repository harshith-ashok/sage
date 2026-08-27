"""Extracts plain text from an uploaded document so app.agent's
read_uploaded_document tool can hand it straight to the reasoning model —
no vision/OCR call needed for already-machine-readable formats. Reuses
pymupdf4llm (already a dependency, Phase 3's PDF ingestion) for PDFs and
python-docx (already a dependency, the .docx *writer* side) for reading
.docx. MS Office reading (.xlsx/.pptx) added alongside the .xlsx/.pptx
*writer* side (app/tasks/excel_writer.py, app/tasks/pptx_writer.py) — same
openpyxl/python-pptx libraries, just the other direction.
"""

import os

import pymupdf
import pymupdf4llm
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation

PLAIN_TEXT_EXTENSIONS = {".txt", ".md", ".csv", ".json", ".yaml", ".yml", ".log"}
SUPPORTED_EXTENSIONS = PLAIN_TEXT_EXTENSIONS | {".pdf", ".docx", ".xlsx", ".pptx"}


class UnsupportedDocumentType(Exception):
    pass


def _extract_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        if rows:
            parts.append("\n".join(rows))
    return "\n\n".join(parts)


def _extract_xlsx(file_path: str) -> str:
    wb = load_workbook(file_path, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = ["\t".join("" if c is None else str(c) for c in row) for row in sheet.iter_rows(values_only=True)]
        rows = [r for r in rows if r.strip("\t")]
        if rows:
            parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()]
        if texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def extract_text(file_path: str, filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        pages = pymupdf4llm.to_markdown(file_path, page_chunks=True)
        return "\n\n".join(f"[Page {i + 1}]\n{page['text']}" for i, page in enumerate(pages))

    if ext == ".docx":
        return _extract_docx(file_path)

    if ext == ".xlsx":
        return _extract_xlsx(file_path)

    if ext == ".pptx":
        return _extract_pptx(file_path)

    if ext in PLAIN_TEXT_EXTENSIONS:
        with open(file_path, encoding="utf-8", errors="replace") as f:
            return f.read()

    raise UnsupportedDocumentType(
        f"Unsupported document type '{ext}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
    )


MAX_RENDERED_PAGES = 5


def render_pdf_pages(file_path: str, max_pages: int = MAX_RENDERED_PAGES, dpi: int = 150) -> list[bytes]:
    """Renders each of a PDF's first `max_pages` pages to a real PNG image
    — `pymupdf4llm.to_markdown()` (extract_text, above) only ever pulls
    *text*, silently dropping any figure/diagram/photo embedded in the
    page; a text-only extraction of a PDF containing a P&ID excerpt or a
    photo would report nothing at all about it. This lets a page also be
    handed to the vision model (app.agent's read_document_page_as_image)
    the same way an uploaded image already is, so a PDF's figures aren't
    invisible to the agent just because they arrived inside a PDF rather
    than as a standalone image."""
    doc = pymupdf.open(file_path)
    try:
        images = []
        zoom = dpi / 72  # PDF's native unit is 72 DPI
        matrix = pymupdf.Matrix(zoom, zoom)
        for page in doc[:max_pages]:
            pix = page.get_pixmap(matrix=matrix)
            images.append(pix.tobytes("png"))
        return images
    finally:
        doc.close()
